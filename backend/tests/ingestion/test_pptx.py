from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app.ingestion.base import BlockType
from app.ingestion.pptx import PptxExtractor


def _build_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text_frame.text = "Net Sales climbed to $6.7M"

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    rows, cols = 2, 2
    table_shape = slide2.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "12345"

    prs.save(path)


def test_extracts_slide_text_and_table(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    _build_pptx(pptx_path)

    blocks = PptxExtractor().extract(pptx_path, document_id="doc-1")

    text_blocks = [b for b in blocks if b.block_type == BlockType.TEXT]
    table_blocks = [b for b in blocks if b.block_type == BlockType.TABLE]

    assert any(b.page == 1 and "Net Sales" in b.text for b in text_blocks)
    assert len(table_blocks) == 1
    assert table_blocks[0].page == 2
    assert table_blocks[0].table == [["Metric", "Value"], ["Revenue", "12345"]]
    assert all(b.document_id == "doc-1" for b in blocks)
