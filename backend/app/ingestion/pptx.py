"""PPTX extractor: slide text and tables via python-pptx. `page` maps to
1-indexed slide number."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.ingestion.base import BlockType, ContentBlock, Extractor


class PptxExtractor(Extractor):
    def extract(self, file_path: Path, document_id: str) -> list[ContentBlock]:
        presentation = Presentation(str(file_path))
        blocks: list[ContentBlock] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []

            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = [[cell.text for cell in row.cells] for row in table.rows]
                    if rows:
                        blocks.append(
                            ContentBlock(
                                document_id=document_id,
                                page=slide_number,
                                block_type=BlockType.TABLE,
                                table=rows,
                            )
                        )
                elif shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text)

            if texts:
                blocks.append(
                    ContentBlock(
                        document_id=document_id,
                        page=slide_number,
                        block_type=BlockType.TEXT,
                        text="\n".join(texts),
                    )
                )

        return blocks
